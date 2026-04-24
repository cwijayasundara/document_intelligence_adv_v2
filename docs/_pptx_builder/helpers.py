"""Reusable shape + text builders for the eval-framework deck."""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .theme import (
    ACCENT,
    BORDER,
    CODE_BG,
    FONT,
    INK,
    MONO,
    MUTED_BG,
    NAVY,
    SLIDE_W,
    SUBTLE,
    TOTAL_SLIDES,
    WHITE,
)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill=WHITE, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=14,
    bold=False,
    color=INK,
    font=FONT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=INK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def header_bar(slide, title, kicker=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=NAVY)
    if kicker:
        add_text(
            slide,
            Inches(0.55),
            Inches(0.30),
            Inches(12),
            Inches(0.32),
            kicker.upper(),
            size=11,
            bold=True,
            color=ACCENT,
        )
        title_y = Inches(0.62)
    else:
        title_y = Inches(0.40)
    add_text(
        slide,
        Inches(0.55),
        title_y,
        Inches(12),
        Inches(0.6),
        title,
        size=28,
        bold=True,
        color=NAVY,
    )
    add_rect(slide, Inches(0.55), Inches(1.30), Inches(0.6), Inches(0.04), fill=ACCENT)


def footer(slide, page_num):
    add_text(
        slide,
        Inches(0.55),
        Inches(7.10),
        Inches(8),
        Inches(0.30),
        "PE Document Intelligence  ·  Evaluation Framework",
        size=9,
        color=SUBTLE,
    )
    add_text(
        slide,
        Inches(11.5),
        Inches(7.10),
        Inches(1.3),
        Inches(0.30),
        f"{page_num} / {TOTAL_SLIDES}",
        size=9,
        color=SUBTLE,
        align=PP_ALIGN.RIGHT,
    )


def card(slide, x, y, w, h, title, lines, *, accent=ACCENT):
    add_rect(slide, x, y, w, h, fill=WHITE, line=BORDER)
    add_rect(slide, x, y, Inches(0.08), h, fill=accent)
    add_text(
        slide,
        x + Inches(0.25),
        y + Inches(0.12),
        w - Inches(0.35),
        Inches(0.42),
        title,
        size=14,
        bold=True,
        color=NAVY,
    )
    body_h = h - Inches(0.62)
    add_bullets(
        slide,
        x + Inches(0.25),
        y + Inches(0.55),
        w - Inches(0.35),
        body_h,
        lines,
        size=11,
        color=INK,
    )


def kpi_tile(slide, x, y, w, h, value, label, *, color=ACCENT):
    add_rect(slide, x, y, w, h, fill=WHITE, line=BORDER)
    add_text(
        slide,
        x,
        y + Inches(0.25),
        w,
        Inches(0.8),
        value,
        size=36,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x,
        y + Inches(1.05),
        w,
        Inches(0.4),
        label,
        size=11,
        color=SUBTLE,
        align=PP_ALIGN.CENTER,
    )


def code_block(slide, x, y, w, h, code, *, size=10):
    add_rect(slide, x, y, w, h, fill=CODE_BG, line=BORDER)
    tb = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.10), w - Inches(0.30), h - Inches(0.20)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = MONO
        run.font.size = Pt(size)
        run.font.color.rgb = INK


def two_col_table(slide, x, y, w, h, rows, col1_w_ratio=0.35):
    n = len(rows)
    row_h = h / n
    col1_w = w * col1_w_ratio
    col2_w = w - col1_w
    for i, (left, right) in enumerate(rows):
        is_header = i == 0
        fill = NAVY if is_header else (MUTED_BG if i % 2 == 0 else WHITE)
        fg_l = WHITE if is_header else NAVY
        fg_r = WHITE if is_header else INK
        add_rect(slide, x, y + row_h * i, col1_w, row_h, fill=fill, line=BORDER)
        add_rect(
            slide, x + col1_w, y + row_h * i, col2_w, row_h, fill=fill, line=BORDER
        )
        add_text(
            slide,
            x + Inches(0.15),
            y + row_h * i,
            col1_w - Inches(0.20),
            row_h,
            left,
            size=11,
            bold=is_header,
            color=fg_l,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            x + col1_w + Inches(0.15),
            y + row_h * i,
            col2_w - Inches(0.20),
            row_h,
            right,
            size=11,
            bold=is_header,
            color=fg_r,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def three_col_table(slide, x, y, w, h, rows, ratios=(0.25, 0.35, 0.40)):
    n = len(rows)
    row_h = h / n
    cw = [w * r for r in ratios]
    cx = [x, x + cw[0], x + cw[0] + cw[1]]
    for i, cells in enumerate(rows):
        is_header = i == 0
        fill = NAVY if is_header else (MUTED_BG if i % 2 == 0 else WHITE)
        fg = WHITE if is_header else INK
        for j, cell in enumerate(cells):
            add_rect(slide, cx[j], y + row_h * i, cw[j], row_h, fill=fill, line=BORDER)
            add_text(
                slide,
                cx[j] + Inches(0.12),
                y + row_h * i,
                cw[j] - Inches(0.18),
                row_h,
                cell,
                size=10.5,
                bold=is_header,
                color=fg,
                anchor=MSO_ANCHOR.MIDDLE,
            )


def make_arrow(slide, x, y, w, h, color=SUBTLE):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow
